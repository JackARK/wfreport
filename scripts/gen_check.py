"""Generate real artifacts for visual/layout verification (not part of pytest).

Usage: uv run python scripts/gen_check.py
Writes to /tmp/wfcheck/{last,current,next}/ and prints a layout report.
"""
import os
import sys
import tempfile
import zipfile
from pathlib import Path

from pptx import Presentation

sys.path.insert(0, str(Path(__file__).resolve().parents[1]))

from backend.core.parser import load_excel
from backend.core.metrics import compute_all
from backend.core import history
from backend.charts.render_png import render_all_pngs
from backend.reports.ppt_builder import build_ppt
from backend.reports.narratives import build_narratives

TEMPLATE = "resource/2026年第二十七周周报-王凡.pptx"
ROOT = Path("/tmp/wfcheck")

PROC = [{"事项内容": f"采购事项{i}", "进度及责任人": ["步骤a", "步骤b"],
         "状态": "进行中" if i % 2 else "已完成", "完成时间": "7月15日"} for i in range(6)]
PLAN = [{"事项内容": f"计划事项{i}", "提出时间": "7月10日",
         "次周预计完成节点名称": "节点", "涉及部门": "供应链"} for i in range(5)]
AI_TEXTS = {"week_compare": "本周销售额较上周上升，毛利保持稳定。（测试文案）",
            "daily_summary": "每日销售平稳，周末略高。（测试文案）"}


def gen_one(xlsx: str, tag: str, db: str):
    df = load_excel(xlsx)
    bundle = compute_all(df)
    week_id = str(df["订单日期"].min().date().isocalendar()[:2])
    y, w = df["订单日期"].min().date().isocalendar()[:2]
    week_id = f"{y}-W{w:02d}"
    history.save_week(week_id, df, bundle.overview, db)
    recent = history.get_recent_weeks(3, db)
    print(f"[{tag}] week_id={week_id} rows={len(df)} new_products={len(bundle.new_products)} "
          f"recent={[r['week_id'] for r in recent]}")

    out_dir = ROOT / tag
    png_dir = out_dir / "png"
    png_dir.mkdir(parents=True, exist_ok=True)
    pngs = render_all_pngs(bundle, recent, str(png_dir), template_path=TEMPLATE)

    week_meta = {"week_id": week_id,
                 "周起始日": str(df["订单日期"].min().date()),
                 "周结束日": str(df["订单日期"].max().date())}
    narratives = build_narratives(bundle, week_meta)
    pptx_path = str(out_dir / f"{week_id}.pptx")
    build_ppt(TEMPLATE, pngs, AI_TEXTS, narratives, PROC, PLAN, week_meta, pptx_path,
              has_new_products=len(bundle.new_products) > 0)
    return pptx_path


def check_layout(pptx_path: str):
    print(f"\n=== layout check: {pptx_path} ===")
    with zipfile.ZipFile(pptx_path) as z:
        names = z.namelist()
        dup = {n for n in names if names.count(n) > 1}
        print(f"zip entries={len(names)} duplicates={sorted(dup) if dup else '无'}")
    prs = Presentation(pptx_path)
    W, H = prs.slide_width, prs.slide_height
    print(f"slide size = {W/914400:.2f}in x {H/914400:.2f}in, slides={len(prs.slides)}")
    problems = []
    for i, s in enumerate(prs.slides):
        title = ""
        for sh in s.shapes:
            if sh.has_text_frame and sh.text_frame.text.strip():
                title = sh.text_frame.text.strip().split("\n")[0][:24]
                break
        pics = sum(1 for sh in s.shapes if sh.shape_type == 13)
        tbls = sum(1 for sh in s.shapes if sh.has_table)
        print(f"[{i:2d}] pics={pics} tbls={tbls} {title}")
        for sh in s.shapes:
            try:
                l, t, w, h = sh.left, sh.top, sh.width, sh.height
            except Exception:
                continue
            if l is None:
                continue
            if l < -9144 or t < -9144 or l + w > W + 9144 or t + h > H + 9144:
                problems.append(f"slide[{i}] shape {sh.name!r} out of bounds: "
                                f"l={l/914400:.2f} t={t/914400:.2f} r={(l+w)/914400:.2f} b={(t+h)/914400:.2f} (in)")
    if problems:
        print("!! 越界形状:")
        for p in problems:
            print("  ", p)
    else:
        print("越界检查: 全部形状在页面内")


def main():
    ROOT.mkdir(exist_ok=True)
    db = os.path.join(tempfile.mkdtemp(), "check.db")
    history.init_db(db)
    p_last = gen_one("resource/上周.xlsx", "last", db)
    p_cur = gen_one("resource/本周.xlsx", "current", db)
    p_next = gen_one("resource/下周.xlsx", "next", db)
    for p in (p_last, p_cur, p_next):
        check_layout(p)


if __name__ == "__main__":
    main()

"""Render plotly figures to PNGs sized to match the PPT picture slots.

Each chart's input data is logged in full (small enough — these are
summary frames, ≤ 15 rows each) so an operator reading logs can
re-derive every figure without re-running the pipeline.
"""
import os
from pathlib import Path
import time
import pandas as pd
from backend.charts import plotly_figures as pf
from backend.core.logging_conf import get_logger

logger = get_logger("render")

# Render resolution knobs. Width in pixels; height is computed from
# the slot's aspect ratio so the image never gets stretched.
_WIDTH_PX = 1400
_DPI = 150
_DEFAULT_W_IN, _DEFAULT_H_IN = 10.0, 5.0  # fallback for charts not in PPT


def _emu_to_in(emu):
    return emu / 914400 if emu is not None else None


def _first_picture(slide):
    for sh in slide.shapes:
        if sh.shape_type == 13:  # PICTURE
            w = _emu_to_in(sh.width)
            h = _emu_to_in(sh.height)
            if w and h:
                return (w, h)
    return None


def _slot_geometry(template_path: str):
    from pptx import Presentation

    if not template_path:
        return {}

    prs = Presentation(template_path)
    S = prs.slides
    geom = {}

    def _pick(slide_idx, key):
        if len(S) <= slide_idx:
            return
        d = _first_picture(S[slide_idx])
        if d:
            geom[key] = d

    if len(S) > 3:
        pics = [sh for sh in S[3].shapes if sh.shape_type == 13]
        if len(pics) >= 2:
            geom["overview"]     = (_emu_to_in(pics[0].width), _emu_to_in(pics[0].height))
            geom["three_weeks"]  = (_emu_to_in(pics[1].width), _emu_to_in(pics[1].height))
        elif len(pics) == 1:
            geom["overview"]     = (_emu_to_in(pics[0].width), _emu_to_in(pics[0].height))

    _pick(4, "daily")
    _pick(5, "brand_combo")
    _pick(6, "brand_pie")
    _pick(7, "shop_heatmap")
    _pick(8, "platform")
    _pick(9, "product_horizontal")
    _pick(10, "factory")
    return geom


def _df_summary(name: str, df: pd.DataFrame | None) -> str:
    """Compact multi-line dump of a dataframe for logs. Caps rows so we
    don't blow up the log file with a 15K-row frame."""
    if df is None or len(df) == 0:
        return f"{name}=<empty>"
    cols = ", ".join(df.columns.tolist())
    head = df.head(5).to_string(index=False, max_cols=8)
    tail = f"...(+{len(df)-5} more)" if len(df) > 5 else ""
    return f"{name} rows={len(df)} cols=[{cols}]\n{head}\n{tail}"


def _scalar_summary(name: str, value) -> str:
    if isinstance(value, dict):
        items = ", ".join(f"{k}={v}" for k, v in value.items())
        return f"{name}={{{items}}}"
    if isinstance(value, list):
        return f"{name}=list[{len(value)}]"
    return f"{name}={value}"


def _log_chart_inputs(name: str, args: dict) -> None:
    """Log the data each chart consumes (bundle.daily, bundle.brand, …)
    so an operator can verify outputs without re-running."""
    parts = [f"chart={name}"]
    for k, v in args.items():
        if isinstance(v, pd.DataFrame):
            parts.append(_df_summary(k, v))
        elif isinstance(v, (dict, list)):
            parts.append(_scalar_summary(k, v))
        else:
            parts.append(f"{k}={v}")
    logger.info("chart inputs\n  %s", "\n  ".join(parts))


def _render_with_aspect(fig, out_path: str, slot_inches=None):
    if slot_inches:
        w_in, h_in = slot_inches
        h_px = max(int(round(_WIDTH_PX * (h_in / w_in))), 200)
    else:
        h_px = int(_WIDTH_PX * (_DEFAULT_H_IN / _DEFAULT_W_IN))
    fig.write_image(out_path, width=_WIDTH_PX, height=h_px, scale=1)


def fig_to_png(fig, out_path: str, width: int = 1200, height: int = 700):
    fig.write_image(out_path, width=width, height=height, scale=1)


def _build_chart_args(name: str, bundle, recent_weeks: list) -> dict:
    """Pull the inputs each chart factory will use, by name. Keeps the
    per-chart logging in one place."""
    if name == "overview":
        return {"bundle.overview": bundle.overview}
    if name == "daily":
        return {"bundle.daily": bundle.daily}
    if name == "brand_combo":
        return {"bundle.brand": bundle.brand}
    if name == "brand_pie":
        return {"bundle.brand": bundle.brand}
    if name == "platform":
        return {"bundle.platform": bundle.platform}
    if name == "shop_heatmap":
        return {"bundle.shop_top15_daily": bundle.shop_top15_daily}
    if name == "product_horizontal":
        return {"bundle.product_top15": bundle.product_top15}
    if name == "product_table":
        return {"bundle.product_top15": bundle.product_top15}
    if name == "product_heatmap":
        return {"bundle.product_top15_daily": bundle.product_top15_daily}
    if name == "new_products":
        return {"bundle.new_products": bundle.new_products}
    if name == "three_weeks":
        return {"recent_weeks": recent_weeks}
    if name == "factory":
        return {"bundle.factory_top5": bundle.factory_top5}
    return {}


def render_all_pngs(bundle, recent_weeks, out_dir: str,
                    template_path: str | None = None) -> dict:
    """Render every chart PNG to `out_dir`. When template_path is given,
    each PNG is sized to match the slot it lands in on the PPT."""
    out = Path(out_dir)
    out.mkdir(parents=True, exist_ok=True)
    geom = _slot_geometry(template_path) if template_path else {}
    logger.info("render_all_pngs START out=%s template=%s slot_dims=%s",
                out_dir, bool(template_path), {k: f"{v[0]:.2f}x{v[1]:.2f}in"
                                                for k, v in geom.items()})

    factories = {
        "overview":           pf.fig_overview,
        "daily":              pf.fig_daily,
        "brand_combo":        pf.fig_brand_combo,
        "brand_pie":          pf.fig_brand_pie,
        "platform":           pf.fig_platform,
        "shop_heatmap":       pf.fig_shop_heatmap,
        "product_horizontal": pf.fig_product_horizontal,
        "product_table":      pf.fig_product_table,
        "product_heatmap":    pf.fig_product_heatmap,
        "new_products":       pf.fig_new_products,
        "three_weeks":        lambda b, r: pf.fig_three_weeks_table(r),
        "factory":            pf.fig_factory_table,
    }

    paths = {}
    t0 = time.perf_counter()
    for name, factory in factories.items():
        chart_t0 = time.perf_counter()
        args = _build_chart_args(name, bundle, recent_weeks)
        _log_chart_inputs(name, args)
        if name == "three_weeks":
            fig = factory(bundle, recent_weeks)
        else:
            fig = factory(bundle)
        p = str(out / f"{name}.png")
        _render_with_aspect(fig, p, geom.get(name))
        size_kb = os.path.getsize(p) / 1024
        slot = geom.get(name)
        slot_str = f"{slot[0]:.2f}x{slot[1]:.2f}in" if slot else "default"
        logger.info("render %s → %s  size=%.0fKB slot=%s 耗时=%.2fs",
                    name, p, size_kb, slot_str, time.perf_counter() - chart_t0)
        paths[name] = p

    logger.info("render_all_pngs DONE charts=%d 总耗时=%.1fs",
                len(paths), time.perf_counter() - t0)
    return paths
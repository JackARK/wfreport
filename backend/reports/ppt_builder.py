import copy
import time
from PIL import Image
from pptx import Presentation

from backend.core.logging_conf import get_logger

logger = get_logger("ppt")

_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
_A_TR = f"{{{_A}}}tr"
_A_T = f"{{{_A}}}t"
_A_TCPR = f"{{{_A}}}tcPr"
_A_B = f"{{{_A}}}b"
_FILL_LOCALS = {"solidFill", "gradFill", "blipFill", "pattFill", "grpFill"}


def _localname(el) -> str:
    tag = el.tag
    return tag.rpartition("}")[2] if isinstance(tag, str) and "}" in tag else (tag or "")


def _strip_row_styling(tr):
    """Remove header-originated bold + cell background from a data <a:tr> clone.

    Strips: (a) any fill child of <a:tcPr> (solidFill/gradFill/...), and
    (b) bold on run props (<a:b/> child or b="1" attribute on rPr/endParaRPr).
    Cell count/structure is left intact.
    """
    for el in tr.iter():
        local = _localname(el)
        if local == "tcPr":
            for child in list(el):
                if _localname(child) in _FILL_LOCALS:
                    el.remove(child)
        elif local in ("rPr", "endParaRPr"):
            for b in el.findall(_A_B):
                el.remove(b)
            if el.get("b") in ("1", "true"):
                del el.attrib["b"]


def _picture_shapes(slide):
    return [s for s in slide.shapes if s.shape_type == 13]  # PICTURE


def _fit_inside(img_w: int, img_h: int, left: int, top: int, width: int, height: int):
    """Contain-fit `img_w x img_h` (px) inside the placeholder box (EMU),
    centered, preserving aspect. Without this, add_picture(width=, height=)
    stretches the chart whenever the PNG aspect != slot aspect."""
    if img_w and img_h and width and height:
        box_ar = width / height
        img_ar = img_w / img_h
        if img_ar > box_ar:
            new_w, new_h = width, int(round(width / img_ar))
        else:
            new_h, new_w = height, int(round(height * img_ar))
        left += (width - new_w) // 2
        top += (height - new_h) // 2
        width, height = new_w, new_h
    return left, top, width, height


def _replace_picture_shape(slide, old, png_path: str):
    """Swap `old` (a Picture shape) in-place: capture its geometry, remove
    the old <p:pic> element from spTree, and insert the new picture
    aspect-fitted inside the same box.

    Also drops the old picture's <Relationship> from the slide's rels so
    the package doesn't carry an orphan reference (image file bloat —
    the embed still references the old PNG even though it's not rendered).
    PowerPoint silently keeps the embed if the rel survives; this is the
    P0-#4 cleanup.
    """
    left, top, width, height = old.left, old.top, old.width, old.height

    # Capture the rId of the old picture's blip embed before we drop the
    # element so we can drop the corresponding relationship.
    blip_rId = None
    blip_ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    embed_attr = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed"
    blip = old._element.find(f".//{{{blip_ns}}}blip")
    if blip is not None:
        blip_rId = blip.get(embed_attr)

    sp = old._element
    sp.getparent().remove(sp)

    # Drop the now-unused relationship + image part from the package.
    if blip_rId:
        try:
            slide.part.drop_rel(blip_rId)
        except Exception:
            pass

    with Image.open(png_path) as im:
        img_w, img_h = im.size
    left, top, width, height = _fit_inside(img_w, img_h, left, top, width, height)
    slide.shapes.add_picture(png_path, left, top, width=width, height=height)


def replace_picture(slide, which: int, png_path: str):
    """Swap the `which`-th picture (document order) with `png_path`.
    NOTE: do NOT call this twice with indices on the same slide — the first
    swap removes a shape and appends the new one at the end of spTree, so
    index 1 afterwards hits the just-inserted picture. Snapshot the shapes
    first and use _replace_picture_shape instead (see build_ppt P4)."""
    pics = _picture_shapes(slide)
    if which >= len(pics):
        return
    _replace_picture_shape(slide, pics[which], png_path)


def set_text(slide, name_part: str, text: str):
    for sh in slide.shapes:
        if sh.has_text_frame and name_part in sh.name:
            tf = sh.text_frame
            tf.clear()
            tf.text = str(text)


def duplicate_slide(prs, index: int):
    src = prs.slides[index]
    blank = prs.slide_layouts[6]
    new = prs.slides.add_slide(blank)
    for sh in src.shapes:
        el = sh._element
        new.shapes._spTree.insert_element_before(copy.deepcopy(el), "p:extLst")
    for rel in src.part.rels.values():
        if "image" in rel.reltype or "chart" in rel.reltype:
            new.part.rels.get_or_add(rel.reltype, rel._target)
    return new


def move_slide(prs, from_idx: int, to_idx: int):
    """Reorder a slide by manipulating the sldIdLst XML directly.
    `to_idx` is the position the slide occupies after the move."""
    sldIdLst = prs.slides._sldIdLst
    elems = list(sldIdLst)
    el = elems[from_idx]
    sldIdLst.remove(el)
    sldIdLst.insert(to_idx, el)


def find_slide_by_title(prs, title: str, *, last: bool = False) -> int | None:
    """Return the index of the slide that has a text shape whose stripped
    text equals `title` exactly. `last=True` picks the last match (skips
    the TOC slide, which repeats section titles)."""
    found = None
    for i, slide in enumerate(prs.slides):
        for sh in slide.shapes:
            if sh.has_text_frame and sh.text_frame.text.strip() == title:
                found = i if last else (found if found is not None else i)
                if not last:
                    return found
    return found


def delete_slide(prs, slide_index: int):
    """Remove a slide from the presentation by index. python-pptx has no
    built-in API for this so we manipulate the sldIdLst XML directly.

    PowerPoint will renumber subsequent slides when reopening. Note that
    slide indices in the caller's list are NOT updated by this call —
    pass the absolute index before any other slides were added/deleted.
    """
    sldIdLst = prs.slides._sldIdLst
    sldId_elems = list(sldIdLst)
    if slide_index < 0 or slide_index >= len(sldId_elems):
        return False
    sldId = sldId_elems[slide_index]
    rId = sldId.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
    # Drop the relationship so the package doesn't carry an orphan
    # reference. Part drop_rel also marks the part for removal.
    try:
        prs.part.drop_rel(rId)
    except Exception:
        pass
    sldIdLst.remove(sldId)
    return True


def _add_row(table, template_tr):
    new_tr = copy.deepcopy(template_tr)
    for t in new_tr.iter(_A_T):
        t.text = ""
    _strip_row_styling(new_tr)
    # append after the current last row so order is preserved
    last_tr = table._tbl.findall(_A_TR)[-1]
    last_tr.addnext(new_tr)


def fill_table(table, items: list, col_keys: list, start_row: int = 1, start_no: int = 1):
    tbl = table._tbl
    trs = tbl.findall(_A_TR)
    # Prefer an existing data-row as the clone template so data rows don't
    # inherit header styling (bold / gradient fill). Fall back to the header
    # row; _strip_row_styling removes header bold/fill from each clone.
    template_tr = trs[start_row] if len(trs) > start_row else trs[0]
    while len(table.rows) > start_row:
        last = len(table.rows) - 1
        tr = table.rows[last]._tr
        tr.getparent().remove(tr)
    for i, item in enumerate(items):
        _add_row(table, template_tr)
        row_idx = len(table.rows) - 1
        table.cell(row_idx, 0).text = str(start_no + i)  # 序号, col 0, continuous
        for c, key in enumerate(col_keys):
            val = item.get(key, "")
            if isinstance(val, list):
                val = "\n".join(f"{j+1}、{x}" for j, x in enumerate(val))
            table.cell(row_idx, c + 1).text = str(val)   # col_keys -> cols 1..N


def build_ppt(template_path, png_paths, ai_texts, narratives, procurement_items,
              plan_items, week_meta, out_path, has_new_products: bool = True) -> str:
    t0 = time.perf_counter()
    logger.info("build_ppt START template=%s out=%s week=%s 采购=%d 计划=%d",
                template_path, out_path, week_meta.get("week_id"),
                len(procurement_items or []), len(plan_items or []))
    logger.info("narratives keys=%s", list(narratives.keys()))
    for k in ("week_compare", "daily_summary"):
        v = ai_texts.get(k, "") or ""
        logger.info("ai_texts.%s len=%d head=%.60s", k, len(v), v[:60].replace("\n", " "))

    prs = Presentation(template_path)
    S = prs.slides
    logger.info("PPT 模板打开 slides=%d", len(prs.slides))

    # 页3 销售表现 — 顶部长条放三周对比表, 右下放每日销售图。
    # 必须先快照两个占位形状再按形状替换: 替换会删除旧图并把新图追加到
    # spTree 末尾, 若按索引连续调用 replace_picture(0/1), 第二次会命中
    # 刚插入的新图 (旧 bug: 三周表被压进右下槽, 长条槽残留模板旧图)。
    set_text(S[3], "文本框 5", ai_texts.get("week_compare", ""))
    p4_pics = _picture_shapes(S[3])
    if len(p4_pics) >= 2:
        _replace_picture_shape(S[3], p4_pics[1], png_paths["three_weeks"])  # 顶部长条
        _replace_picture_shape(S[3], p4_pics[0], png_paths["daily"])        # 右下
    elif len(p4_pics) == 1:
        _replace_picture_shape(S[3], p4_pics[0], png_paths["three_weeks"])
    # 页4 每日
    replace_picture(S[4], 0, png_paths["daily"])
    # 页5 品牌
    set_text(S[5], "文本框 9", narratives.get("brand", ""))
    replace_picture(S[5], 0, png_paths["brand_combo"])
    # 页6 品牌饼
    set_text(S[6], "文本框 2", narratives.get("brand_share", ""))
    replace_picture(S[6], 0, png_paths["brand_pie"])
    # 页7 店铺热力图
    replace_picture(S[7], 0, png_paths["shop_heatmap"])
    # 页8 平台
    replace_picture(S[8], 0, png_paths["platform"])
    # 页9 产品
    set_text(S[9], "文本框 5", narratives.get("product", ""))
    replace_picture(S[9], 0, png_paths["product_horizontal"])
    # 页10 工厂 (供应商 TOP5)
    replace_picture(S[10], 0, png_paths["factory"])
    # 采购自动分页 (页13,14,15 in 1-indexed -> [12,13,14] 0-indexed)
    # - Add pages if items > 12; inserted right after the procurement block
    #   so they never land behind the THANKS page.
    # - DELETE unused template pages when items < 12, otherwise the leftover
    #   template content ("京东自营及猫超产品入库" etc.) ships in the PPT.
    # - 0 items: keep exactly 1 page filled with an empty (header-only) table
    #   so the 采购情况 section never ships with zero content pages.
    max_per = 4
    chunks = [procurement_items[i:i + max_per] for i in range(0, len(procurement_items), max_per)]
    proc_pages = [12, 13, 14]
    while len(chunks) > len(proc_pages):
        duplicate_slide(prs, proc_pages[-1])
        move_slide(prs, len(prs.slides) - 1, proc_pages[-1] + 1)
        proc_pages.append(proc_pages[-1] + 1)
    keep = max(len(chunks), 1)
    logger.info("采购自动分页 chunks=%d proc_pages=%s keep=%d", len(chunks), proc_pages, keep)
    running_no = 1
    for pi in range(keep):
        chunk = chunks[pi] if pi < len(chunks) else []
        slide = prs.slides[proc_pages[pi]]
        for sh in slide.shapes:
            if sh.has_table:
                logger.info("填充采购表 slide=%d rows=%d start_no=%d",
                            proc_pages[pi] + 1, len(chunk), running_no)
                fill_table(sh.table, chunk, ["事项内容", "进度及责任人", "状态", "完成时间"], start_no=running_no)
        running_no += len(chunk)
    # 新品页 — 有新品时复制页9 放新品图，稍后插入到「下周规划」分隔页之前；
    # 无新品时整页不生成。复制必须在删除多余采购页之前做：python-pptx 的
    # add_slide 按 len(slides)+1 分配 partname，先删后加会与现存 slide 的
    # partname 撞名，保存出 duplicate zip 条目。
    new9 = None
    if has_new_products:
        new9 = duplicate_slide(prs, 9)
        for sh in new9.shapes:
            if sh.has_text_frame and "文本框 5" in sh.name:
                sh.text_frame.text = narratives.get("new", "")
            elif sh.has_text_frame and sh.text_frame.text.strip() == "销售表现-产品分析":
                sh.text_frame.text = "销售表现-新品分析"
        # P0-#4: drop the inherited product chart's rel (and any other
        # picture rels the duplicate brought along) before adding the
        # new_products chart. Use replace_picture so the geometry + rel
        # cleanup logic is consistent with the main picture swap.
        if _picture_shapes(new9):
            replace_picture(new9, 0, png_paths["new_products"])
    else:
        logger.info("无新品 — 跳过新品页")

    # P0-#3: drop unused procurement slides. We must do this AFTER the
    # fill_table loop so all the slides we're keeping have their data in
    # place. Delete in DESCENDING index order so earlier deletions don't
    # shift the indexes of later ones.
    for unused_idx in sorted([proc_pages[i] for i in range(keep, len(proc_pages))],
                             reverse=True):
        delete_slide(prs, unused_idx)
        logger.info("删除多余采购 slide idx=%d", unused_idx)

    if new9 is not None:
        # 删除采购页后 new9 的索引已前移，按 slide_id 重新定位再移动。
        new9_idx = next(i for i, s in enumerate(prs.slides) if s.slide_id == new9.slide_id)
        plan_divider_idx = find_slide_by_title(prs, "下周规划", last=True)
        if plan_divider_idx is not None:
            move_slide(prs, new9_idx, plan_divider_idx)
            logger.info("新品页插入到 slide idx=%d (下周规划 之前)", plan_divider_idx)
        else:
            logger.warning("未找到「下周规划」分隔页 — 新品页保留在末尾")

    # 下周计划 — locate the slide by content rather than by index. The
    # procurement-slide deletion above shifts indexes by -1, which used
    # to silently land us on the "THANKS" / 感谢观看 page when only
    # 1-2 procurement pages were filled.
    plan_keys = ["事项内容", "提出时间", "次周预计完成节点名称", "涉及部门"]
    plan_filled = False
    for slide in prs.slides:
        for sh in slide.shapes:
            if not sh.has_table:
                continue
            t = sh.table
            if len(t.columns) < len(plan_keys):
                continue
            header = [c.text_frame.text.strip() for c in t.rows[0].cells]
            if all(any(k in h for h in header) for k in plan_keys):
                logger.info("填充下周计划表 rows=%d", len(plan_items or []))
                fill_table(t, plan_items, plan_keys, start_no=1)
                plan_filled = True
                break
        else:
            continue
        break
    if not plan_filled:
        logger.warning("未找到下周计划表格 — 模板中可能缺少匹配表头")

    prs.save(out_path)
    import os as _os
    logger.info("build_ppt DONE → %s size=%.0fKB 总耗时=%.1fs",
                out_path, _os.path.getsize(out_path) / 1024, time.perf_counter() - t0)
    return out_path

import os, zipfile
from pathlib import Path
from openpyxl import load_workbook
from pptx import Presentation
from backend.core.parser import load_excel
from backend.core.metrics import compute_all
from backend.charts.render_png import render_all_pngs
from backend.reports.excel_builder import build_excel
from backend.reports.ppt_builder import build_ppt


def test_full_pipeline(tmp_path):
    df = load_excel("resource/本周.xlsx")
    bundle = compute_all(df)
    recent = []  # single week; 近三周 just empty here
    png_dir = tmp_path / "png"
    pngs = render_all_pngs(bundle, recent, str(png_dir))
    assert set(pngs.keys()) == {"overview", "daily", "brand_combo", "brand_pie", "platform",
                                "shop_heatmap", "product_horizontal", "product_table",
                                "product_heatmap", "new_products", "three_weeks", "factory"}
    # Excel
    xlsx = tmp_path / "out.xlsx"
    build_excel(bundle, recent, {"week_compare": "测试", "daily_summary": "", "procurement": "", "next_plan": ""},
                [], [], str(xlsx))
    wb = load_workbook(str(xlsx))
    assert {"汇总", "每日", "品牌", "平台", "店铺", "产品", "工厂", "AI文案"} <= set(wb.sheetnames)
    with zipfile.ZipFile(str(xlsx)) as z:
        assert any("xl/charts/chart" in n for n in z.namelist())
    # PPT
    pptx = tmp_path / "out.pptx"
    proc = [{"事项内容": "测试采购", "进度及责任人": ["a", "b"], "状态": "已完成", "完成时间": "7月9日"}] * 6
    plan = [{"事项内容": "测试计划", "提出时间": "7月10日", "次周预计完成节点名称": "节点", "涉及部门": "供应链"}] * 5
    narratives = {"brand": "B", "brand_share": "S", "product": "P", "new": "N"}
    build_ppt("resource/2026年第二十七周周报-王凡.pptx", pngs,
              {"week_compare": "测试", "daily_summary": ""}, narratives, proc, plan,
              {"week_id": "2026-W27", "周起始日": "2026-07-03", "周结束日": "2026-07-09"}, str(pptx),
              has_new_products=len(bundle.new_products) > 0)
    prs = Presentation(str(pptx))
    # 原始18 + 新品页1 - 未用采购页1 (6条采购=2页, 模板3页删1)
    assert len(prs.slides) == 18
    titles = []
    for s in prs.slides:
        t = ""
        for sh in s.shapes:
            if sh.has_text_frame and sh.text_frame.text.strip():
                t = sh.text_frame.text.strip().split("\n")[0]
                break
        titles.append(t)
    new_idx = titles.index("销售表现-新品分析")
    plan_divider_idx = max(i for i, t in enumerate(titles) if t == "下周规划")
    assert new_idx == plan_divider_idx - 1

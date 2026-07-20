# tests/reports/test_ppt_builder.py
import os, tempfile, base64
from pptx import Presentation
from backend.reports.ppt_builder import build_ppt

_VALID_PNG = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mNk+M9QDwADhgGAWjR9awAAAABJRU5ErkJggg=="
)

_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
_A_TR = f"{{{_A}}}tr"
_A_TCPR = f"{{{_A}}}tcPr"
_A_RPR = f"{{{_A}}}rPr"
_A_ENDPARARPR = f"{{{_A}}}endParaRPr"
_FILL_LOCALS = {"solidFill", "gradFill", "blipFill", "pattFill", "grpFill"}


def _pngs(d):
    paths = {}
    for name in ["overview", "daily", "brand_combo", "brand_pie", "platform",
                 "shop_heatmap", "product_horizontal", "product_table",
                 "product_heatmap", "new_products", "three_weeks", "factory"]:
        p = os.path.join(d, name + ".png")
        with open(p, "wb") as f:
            f.write(_VALID_PNG)
        paths[name] = p
    return paths


def _slide_titles(prs):
    """One-line title text per slide (first non-empty text shape)."""
    titles = []
    for s in prs.slides:
        t = ""
        for sh in s.shapes:
            if sh.has_text_frame and sh.text_frame.text.strip():
                t = sh.text_frame.text.strip().split("\n")[0]
                break
        titles.append(t)
    return titles


def _localname(el):
    tag = el.tag
    return tag.rpartition("}")[2] if isinstance(tag, str) and "}" in tag else (tag or "")


def test_build_ppt_smoke():
    d = tempfile.mkdtemp()
    out = build_ppt(
        template_path="resource/2026年第二十七周周报-王凡.pptx",
        png_paths=_pngs(d),
        ai_texts={"week_compare": "本周概况测试", "daily_summary": "每日测试"},
        narratives={"brand": "品牌叙述", "product": "产品叙述", "new": "新品叙述"},
        procurement_items=[{"事项内容": "测试采购", "进度及责任人": ["a", "b"],
                            "状态": "已完成", "完成时间": "7月9日"}] * 6,
        plan_items=[{"事项内容": "测试计划", "提出时间": "7月10日",
                     "次周预计完成节点名称": "节点", "涉及部门": "供应链"}] * 5,
        week_meta={"week_id": "2026-W27", "周起始日": "2026-07-03", "周结束日": "2026-07-09"},
        out_path=os.path.join(d, "r.pptx"),
    )
    assert os.path.exists(out)
    prs = Presentation(out)
    # 原始18 + 新品页1 - 未用采购页1 (6条采购=2页, 模板3页删1)
    assert len(prs.slides) == 18

    # 验证采购表(页12)列对齐: col0=序号, col1=事项内容, col4=完成时间
    table = None
    for sh in prs.slides[12].shapes:
        if sh.has_table:
            table = sh.table
            break
    assert table is not None, "页12 未找到表格"
    assert table.cell(1, 0).text == "1", f"序号列错误: {table.cell(1, 0).text!r}"
    assert table.cell(1, 1).text == "测试采购", f"事项内容列错误: {table.cell(1, 1).text!r}"
    assert table.cell(1, 4).text == "7月9日", f"完成时间列错误: {table.cell(1, 4).text!r}"


def test_new_products_page_before_next_week_plan():
    """有新品: 新品页(销售表现-新品分析)必须插在「下周规划」分隔页之前。"""
    d = tempfile.mkdtemp()
    out = build_ppt(
        template_path="resource/2026年第二十七周周报-王凡.pptx",
        png_paths=_pngs(d),
        ai_texts={"week_compare": "", "daily_summary": ""},
        narratives={"brand": "", "brand_share": "", "product": "", "new": "新品叙述"},
        procurement_items=[], plan_items=[],
        week_meta={"week_id": "2026-W27", "周起始日": "2026-07-03", "周结束日": "2026-07-09"},
        out_path=os.path.join(d, "order.pptx"),
        has_new_products=True,
    )
    prs = Presentation(out)
    titles = _slide_titles(prs)
    new_idx = titles.index("销售表现-新品分析")
    plan_divider_idx = max(i for i, t in enumerate(titles) if t == "下周规划")
    assert new_idx == plan_divider_idx - 1, (
        f"新品页 idx={new_idx} 未紧邻 下周规划分隔页 idx={plan_divider_idx}: {titles}")
    assert titles[-1].startswith("THANKS") or "感谢" in titles[-1], titles[-1]


def test_no_new_products_skips_new_page():
    """无新品: 不新增新品页, 成品中不存在「销售表现-新品分析」页。"""
    d = tempfile.mkdtemp()
    out = build_ppt(
        template_path="resource/2026年第二十七周周报-王凡.pptx",
        png_paths=_pngs(d),
        ai_texts={"week_compare": "", "daily_summary": ""},
        narratives={"brand": "", "brand_share": "", "product": "", "new": "本周无新品。"},
        procurement_items=[{"事项内容": "x", "进度及责任人": ["a"],
                            "状态": "s", "完成时间": "t"}] * 6,
        plan_items=[],
        week_meta={"week_id": "2026-W27", "周起始日": "2026-07-03", "周结束日": "2026-07-09"},
        out_path=os.path.join(d, "nonew.pptx"),
        has_new_products=False,
    )
    prs = Presentation(out)
    titles = _slide_titles(prs)
    assert "销售表现-新品分析" not in titles
    # 原始18 - 未用采购页1 (6条=2页删1), 无新品页
    assert len(prs.slides) == 17


def test_zero_procurement_keeps_one_empty_page():
    """0 条采购: 保留 1 页仅有表头的采购表, 不出现模板残留内容, 也不整节删光。"""
    d = tempfile.mkdtemp()
    out = build_ppt(
        template_path="resource/2026年第二十七周周报-王凡.pptx",
        png_paths=_pngs(d),
        ai_texts={"week_compare": "", "daily_summary": ""},
        narratives={"brand": "", "brand_share": "", "product": "", "new": ""},
        procurement_items=[], plan_items=[],
        week_meta={"week_id": "2026-W27", "周起始日": "2026-07-03", "周结束日": "2026-07-09"},
        out_path=os.path.join(d, "zeroproc.pptx"),
        has_new_products=False,
    )
    prs = Presentation(out)
    proc_tables = []
    for s in prs.slides:
        for sh in s.shapes:
            if sh.has_table:
                header = [c.text_frame.text.strip() for c in sh.table.rows[0].cells]
                if any("进度及责任人" in h for h in header):
                    proc_tables.append(sh.table)
    assert len(proc_tables) == 1, f"采购表页数={len(proc_tables)}, 期望 1"
    assert len(proc_tables[0].rows) == 1, "0 条采购时表应只剩表头"
    body_text = ""
    for s in prs.slides:
        for sh in s.shapes:
            if sh.has_text_frame:
                body_text += sh.text_frame.text
            if sh.has_table:
                for row in sh.table.rows:
                    for c in row.cells:
                        body_text += c.text_frame.text
    assert "京东自营及猫超产品入库" not in body_text, "模板残留内容未清除"


def test_extra_procurement_pages_stay_before_next_week_plan():
    """采购超过12条时增页必须紧跟采购区, 不能跑到 THANKS 之后。"""
    d = tempfile.mkdtemp()
    items = [{"事项内容": f"采购{i}", "进度及责任人": ["a"],
              "状态": "s", "完成时间": "t"} for i in range(14)]
    out = build_ppt(
        template_path="resource/2026年第二十七周周报-王凡.pptx",
        png_paths=_pngs(d),
        ai_texts={"week_compare": "", "daily_summary": ""},
        narratives={"brand": "", "brand_share": "", "product": "", "new": "新品叙述"},
        procurement_items=items, plan_items=[],
        week_meta={"week_id": "2026-W27", "周起始日": "2026-07-03", "周结束日": "2026-07-09"},
        out_path=os.path.join(d, "proc.pptx"),
    )
    prs = Presentation(out)
    titles = _slide_titles(prs)
    # 14条采购 → 4页 (模板3页 + 增1页); 序号连续 1..14
    proc_rows = []
    for s in prs.slides:
        for sh in s.shapes:
            if sh.has_table:
                t = sh.table
                header = [c.text_frame.text.strip() for c in t.rows[0].cells]
                if any("进度及责任人" in h for h in header):
                    for r in range(1, len(t.rows)):
                        proc_rows.append(t.cell(r, 0).text)
    assert proc_rows == [str(i) for i in range(1, 15)], proc_rows
    # 所有采购表页都在 下周规划 之前, THANKS 在最后
    plan_divider_idx = max(i for i, t in enumerate(titles) if t == "下周规划")
    new_idx = titles.index("销售表现-新品分析")
    assert new_idx == plan_divider_idx - 1
    assert "感谢" in titles[-1] or titles[-1].startswith("THANKS")


def test_data_rows_do_not_inherit_header_styling():
    """Data rows must not carry the header's bold/background fill."""
    d = tempfile.mkdtemp()
    out = build_ppt(
        template_path="resource/2026年第二十七周周报-王凡.pptx",
        png_paths=_pngs(d),
        ai_texts={"week_compare": "", "daily_summary": ""},
        narratives={"brand": "", "brand_share": "", "product": "", "new": ""},
        procurement_items=[{"事项内容": "x", "进度及责任人": ["a"],
                            "状态": "s", "完成时间": "t"}] * 2,
        plan_items=[],
        week_meta={"week_id": "2026-W27", "周起始日": "2026-07-03", "周结束日": "2026-07-09"},
        out_path=os.path.join(d, "styling.pptx"),
    )
    prs = Presentation(out)
    table = None
    for sh in prs.slides[12].shapes:
        if sh.has_table:
            table = sh.table
            break
    assert table is not None
    trs = table._tbl.findall(_A_TR)
    assert len(trs) >= 3, f"expected header + 2 data rows, got {len(trs)}"
    # inspect every data row (skip header at index 0)
    for tr in trs[1:]:
        # 1) no background fill on any cell's tcPr
        for tcpr in tr.iter(_A_TCPR):
            for child in tcpr:
                assert _localname(child) not in _FILL_LOCALS, (
                    f"data row carries header cell fill: {_localname(child)}")
        # 2) no bold runs (neither b="1" attr nor <a:b/> child)
        for rpr in tr.iter(_A_RPR):
            assert rpr.get("b") not in ("1", "true"), "data-row run is bold (b attr)"
            assert len(rpr.findall(f"{{{_A}}}b")) == 0, "data-row run is bold (<a:b/>)"
        for epr in tr.iter(_A_ENDPARARPR):
            assert epr.get("b") not in ("1", "true"), "data-row endParaRPr is bold"


def test_data_row_styling_when_table_has_only_header():
    """Edge case: a table reduced to header-only still clones without header styling."""
    d = tempfile.mkdtemp()
    # next_plan table on slide 16 has 10 rows normally; fill with 1 item to
    # force the clone path. We assert the resulting data row has no grad fill.
    out = build_ppt(
        template_path="resource/2026年第二十七周周报-王凡.pptx",
        png_paths=_pngs(d),
        ai_texts={"week_compare": "", "daily_summary": ""},
        narratives={"brand": "", "brand_share": "", "product": "", "new": ""},
        procurement_items=[],
        plan_items=[{"事项内容": "y", "提出时间": "t", "次周预计完成节点名称": "n",
                     "涉及部门": "d"}],
        week_meta={"week_id": "2026-W27", "周起始日": "2026-07-03", "周结束日": "2026-07-09"},
        out_path=os.path.join(d, "header_only.pptx"),
    )
    prs = Presentation(out)
    table = None
    for slide in prs.slides:
        for sh in slide.shapes:
            if not sh.has_table:
                continue
            header = [c.text_frame.text.strip() for c in sh.table.rows[0].cells]
            if any("次周预计完成节点名称" in h for h in header):
                table = sh.table
                break
        if table is not None:
            break
    assert table is not None, "下周计划表未找到"
    trs = table._tbl.findall(_A_TR)
    assert len(trs) >= 2
    # the single data row (index 1) must not carry the header gradient fill
    for tcpr in trs[1].iter(_A_TCPR):
        for child in tcpr:
            assert _localname(child) not in _FILL_LOCALS, (
                f"data row carries header fill: {_localname(child)}")


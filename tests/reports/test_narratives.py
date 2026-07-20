# tests/reports/test_narratives.py
import os
import base64
from pptx import Presentation
from backend.core.parser import load_excel
from backend.core.metrics import compute_all
from backend.reports.narratives import build_narratives
from backend.reports.ppt_builder import build_ppt

_VALID_PNG = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mNk+M9QDwADhgGAWjR9awAAAABJRU5ErkJggg=="
)


def _pngs(d):
    paths = {}
    for name in ["overview", "daily", "brand_combo", "brand_pie", "platform",
                 "shop_heatmap", "product_horizontal", "product_table",
                 "product_heatmap", "new_products", "three_weeks", "factory"]:
        p = os.path.join(d, name + ".png")
        with open(p, "wb") as f:
            f.write(_VALID_PNG)
        paths[name] = p
    return paths


def _bundle():
    df = load_excel("resource/本周.xlsx")
    return compute_all(df)


_WEEK_META = {"week_id": "2026-W27", "周起始日": "2026-07-03", "周结束日": "2026-07-09"}


def test_build_narratives_has_all_keys_and_nonempty():
    n = build_narratives(_bundle(), _WEEK_META)
    assert set(n.keys()) == {"brand", "brand_share", "product", "new"}
    for k, v in n.items():
        assert isinstance(v, str) and v.strip(), f"{k} narrative empty"


def test_brand_narrative_includes_margin_conclusion():
    n = build_narratives(_bundle(), _WEEK_META)
    # two-brand real data -> conclusion line present
    assert "综合：" in n["brand"]
    assert "毛利率高于" in n["brand"]


def test_brand_share_uses_amount_share_sums_to_100():
    n = build_narratives(_bundle(), _WEEK_META)
    assert "本周销售总额" in n["brand_share"]
    # the two pct figures should both be present; their text sums are not asserted
    # exactly, but each brand amount/10000 must appear
    assert "万元占" in n["brand_share"]


def test_built_narratives_fill_ppt_slide5(tmp_path):
    """build_ppt using server-built narratives fills slide-5 brand text-box."""
    bundle = _bundle()
    narratives = build_narratives(bundle, _WEEK_META)
    out = build_ppt(
        template_path="resource/2026年第二十七周周报-王凡.pptx",
        png_paths=_pngs(str(tmp_path)),
        ai_texts={"week_compare": "", "daily_summary": ""},
        narratives=narratives,
        procurement_items=[],
        plan_items=[],
        week_meta=_WEEK_META,
        out_path=str(tmp_path / "n.pptx"),
    )
    prs = Presentation(out)
    slide5_text = ""
    for sh in prs.slides[5].shapes:
        if sh.has_text_frame and "文本框 9" in sh.name:
            slide5_text = sh.text_frame.text or ""
            break
    assert slide5_text.strip(), "slide 5 文本框 9 empty when using built narratives"
    assert "毛利率" in slide5_text

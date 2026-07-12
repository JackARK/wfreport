# tests/reports/test_ppt_shape_names.py
# Verifies the actual shape.name values on slides 3/5/6/9 (and duplicated 9')
# match what ppt_builder.set_text targets, and that the provided text actually
# lands in the target text-box (catches silent no-op mismatches).
import os
import tempfile
import base64
from pptx import Presentation
from backend.reports.ppt_builder import build_ppt

_VALID_PNG = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mNk+M9QDwADhgGAWjR9awAAAABJRU5ErkJggg=="
)

# (slide_index, expected_target_shape_name_substring) per ppt_builder.build_ppt
_EXPECTED = {
    3: "文本框 5",   # set_text(S[3], "文本框 5", ai_texts["week_compare"])
    5: "文本框 9",   # set_text(S[5], "文本框 9", narratives["brand"])
    6: "文本框 2",   # set_text(S[6], "文本框 2", narratives["brand_share"])
    9: "文本框 5",   # set_text(S[9], "文本框 5", narratives["product"])
}


def _pngs(d):
    paths = {}
    for name in ["overview", "daily", "brand_combo", "brand_pie", "platform",
                 "shop_heatmap", "product_combo", "product_heatmap", "new_products",
                 "three_weeks", "factory"]:
        p = os.path.join(d, name + ".png")
        with open(p, "wb") as f:
            f.write(_VALID_PNG)
        paths[name] = p
    return paths


def _find_text(slide, name_part):
    """Return concatenated text of the first text shape whose name contains name_part."""
    for sh in slide.shapes:
        if sh.has_text_frame and name_part in sh.name:
            return sh.text_frame.text or ""
    return None


def _build(d):
    ai_texts = {"week_compare": "WC_SENTINEL_上周对比", "daily_summary": "每日测试"}
    narratives = {
        "brand": "BRAND_SENTINEL_品牌叙述",
        "brand_share": "SHARE_SENTINEL_占比叙述",
        "product": "PRODUCT_SENTINEL_产品叙述",
        "new": "NEW_SENTINEL_新品叙述",
    }
    out = os.path.join(d, "shape_names.pptx")
    build_ppt(
        template_path="resource/2026年第二十七周周报-王凡.pptx",
        png_paths=_pngs(d),
        ai_texts=ai_texts,
        narratives=narratives,
        procurement_items=[],
        plan_items=[],
        week_meta={"week_id": "2026-W27", "周起始日": "2026-07-03", "周结束日": "2026-07-09"},
        out_path=out,
    )
    return out


def test_target_shape_names_exist_on_template():
    """The shape names ppt_builder targets must exist verbatim on the template slides."""
    prs = Presentation("resource/2026年第二十七周周报-王凡.pptx")
    for idx, name_part in _EXPECTED.items():
        names = [sh.name for sh in prs.slides[idx].shapes if sh.has_text_frame]
        assert any(name_part == n for n in names), (
            f"slide {idx}: target shape {name_part!r} not found among text shapes {names}")


def test_built_ppt_target_textboxes_contain_sentinel_text():
    """Distinct sentinel text must land in the target text-box on each slide."""
    d = tempfile.mkdtemp()
    out = _build(d)
    prs = Presentation(out)
    sentinels = {
        3: "WC_SENTINEL_上周对比",
        5: "BRAND_SENTINEL_品牌叙述",
        6: "SHARE_SENTINEL_占比叙述",
        9: "PRODUCT_SENTINEL_产品叙述",
    }
    for idx, sentinel in sentinels.items():
        name_part = _EXPECTED[idx]
        txt = _find_text(prs.slides[idx], name_part)
        assert txt is not None, f"slide {idx}: target shape {name_part!r} missing on built PPT"
        assert sentinel in txt, (
            f"slide {idx}: sentinel {sentinel!r} not in target textbox {name_part!r}; got {txt!r}")


def test_duplicated_slide9_prime_contains_new_narrative():
    """The duplicated 9' slide carries the 'new' narrative in its copied 文本框 5."""
    d = tempfile.mkdtemp()
    out = _build(d)
    prs = Presentation(out)
    # 9' is appended after the original 18 slides; find the last slide whose
    # 文本框 5 text-box contains the NEW sentinel.
    found = False
    for slide in prs.slides:
        txt = _find_text(slide, "文本框 5")
        if txt and "NEW_SENTINEL_新品叙述" in txt:
            found = True
            break
    assert found, "duplicated 9' slide did not receive the 'new' narrative"
